# app.py

import streamlit as st
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re
from datetime import datetime, timedelta

# ──────────────────────────────────────────────────────────────────────────────
#  Helper: Fill PPTX template with invoice data (finding tables by alt_text)
# ──────────────────────────────────────────────────────────────────────────────
def generate_filled_invoice(rows, template_path, bill_info, payment_method):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Fill text boxes (titles + values)
    text_fields = {
        "Bill No": bill_info.get("Bill No", ""),
        "Bill Date": bill_info.get("Bill Date", ""),
        "Due Date": bill_info.get("Due Date", ""),
        "Biller Name": bill_info.get("Biller Name", ""),
        "Client Address": bill_info.get("Client Address", ""),
        "Client Phone Number": bill_info.get("Client Phone Number", ""),
        "Client Email": bill_info.get("Client Email", ""),
        "Client Bill To": bill_info.get("Client Bill To", ""),
    }
    field_titles = {
        "Bill No": "Bill No: ",
        "Bill Date": "Bill Date: ",
        "Due Date": "Due Date: ",
        "Biller Name": "Biller Name: ",
        "Client Bill To": "Bill To: ",
        "Client Address": "Address: ",
        "Client Phone Number": "Phone: ",
        "Client Email": "Email ID: ",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            name = getattr(shape, "name", "")
            if name in text_fields and name in field_titles:
                value = str(text_fields[name])
                if name == "Client Address":
                    value = value[:65]
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                run_title = p.add_run()
                run_title.text = field_titles[name]
                run_title.font.bold = True
                run_title.font.name = "Poppins"
                run_title.font.size = Pt(12)
                run_value = p.add_run()
                run_value.text = value
                run_value.font.bold = False
                run_value.font.name = "Poppins"
                run_value.font.size = Pt(12)
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Handle Payment Method Checkboxes
    checkbox_names = {
        "Cash": "Cash Check",
        "NEFT / IMPS": "NEFT Check",
        "UPI": "UPI Check",
        "Cheque": "Cheque Check",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            name = getattr(shape, "name", "")
            if name in checkbox_names.values():
                if name == checkbox_names.get(payment_method):
                    shape.text = "✔"
                    # Center align horizontally and vertically
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    # Optionally set font for tick
                    for run in shape.text_frame.paragraphs[0].runs:
                        run.font.name = "Poppins"
                        run.font.size = Pt(12)  # Adjust size as needed
                else:
                    shape.text = ""

    # Fill tables (LineItems and BillingSummary)
    line_table = None
    summary_table = None
    for shape in slide.shapes:
        if shape.has_table:
            if getattr(shape, "name", "") == "LineItems":
                line_table = shape.table
            elif getattr(shape, "name", "") == "BillingSummary":
                summary_table = shape.table

    if line_table is None or summary_table is None:
        raise Exception("Could not find LineItems or BillingSummary table in template.")

    # Get font style from first data cell
    max_rows = len(line_table.rows)
    if max_rows > 1:
        sample_cell = line_table.rows[1].cells[0]
    else:
        sample_cell = line_table.rows[0].cells[0]
    sample_para = sample_cell.text_frame.paragraphs[0]
    sample_run = sample_para.runs[0] if sample_para.runs else None
    if sample_run:
        base_font_name = sample_run.font.name or "Poppins"
        base_font_size = sample_run.font.size or Pt(12)
    else:
        base_font_name = "Poppins"
        base_font_size = Pt(12)

    # Clear existing data rows (rows 1…end), leave row 0 intact
    available_data_rows = max_rows - 1
    for r_idx in range(1, max_rows):
        for c_idx in range(len(line_table.columns)):
            cell = line_table.rows[r_idx].cells[c_idx]
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER

    # Write each data row into rows 1…up to available_data_rows
    rows_to_fill = min(len(rows), available_data_rows)
    for i in range(rows_to_fill):
        row_data = rows[i]
        target_row = line_table.rows[i + 1]  # +1 to skip header
        for col_idx, key in enumerate(
            ["No.", "Item Description", "Weight", "Rate (₹)", "Amount (₹)"]
        ):
            value = row_data.get(key, "")
            if key == "Amount (₹)":
                try:
                    amount_val = float(value)
                    txt = f"{amount_val:,.2f}"
                except Exception:
                    txt = str(value)
            else:
                txt = str(value) if value is not None else ""
            cell = target_row.cells[col_idx]
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = txt
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER

    # Compute Subtotal, Rounding, Net Payable
    amounts = []
    for row_data in rows:
        amt = pd.to_numeric(row_data.get("Amount (₹)", 0), errors="coerce")
        if pd.isna(amt):
            amt = 0.0
        amounts.append(float(amt))
    subtotal = sum(amounts)
    rounded_total = float(round(subtotal))
    rounding_value = rounded_total - subtotal
    net_payable = rounded_total

    # Fill BillingSummary table (row 1: Subtotal, row 2: Rounding, row 3: NET PAYABLE)
    def set_summary_cell(r_idx, value, prefix=""):
        cell = summary_table.rows[r_idx].cells[1]
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = f"{prefix}{value:,.2f}"
        run.font.name = base_font_name
        run.font.size = base_font_size
        para.alignment = PP_ALIGN.CENTER

    set_summary_cell(1, subtotal)
    set_summary_cell(2, rounding_value)
    set_summary_cell(3, net_payable, prefix="₹ ")

    # Save PPTX to memory
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
#  Streamlit App
# ──────────────────────────────────────────────────────────────────────────────
def parse_number(val):
    """
    Parse a number from a string, stripping spaces, commas, Rs, ₹, etc.
    Returns float or raises ValueError.
    """
    if val is None:
        raise ValueError("Empty value")
    s = str(val)
    s = s.replace(",", "")
    s = re.sub(r"(rs\.?|₹)", "", s, flags=re.IGNORECASE)
    s = s.strip()
    if not s:
        raise ValueError("Empty value")
    return float(s)

def main():
    st.set_page_config(
        page_title="💎 Rishab Gems ‒ Diamond Jewellery Invoice Generator",
        layout="wide",
    )

    st.title("💎 Rishab Gems ‒ Diamond Jewellery Invoice Generator")
    st.markdown(
        """
        Fill in your invoice line items below. Press “➕ Add Another Row” to append more rows.
        Once you’re done, click **Generate Invoice** to download a fully-filled PowerPoint (.pptx).
        """
    )

    # ── NEW: Bill Info Section ──
    today = datetime.today().date()
    default_due = today + timedelta(days=7)

    # Generate Bill No: e.g. RG-20240601-001
    if "bill_no" not in st.session_state:
        st.session_state.bill_no = f"RG-{today.strftime('%Y%m%d')}-{datetime.now().strftime('%H%M%S')}"

    bill_no = st.session_state.bill_no
    bill_date = st.date_input("Bill Date", value=today, key="bill_date")
    due_date = st.date_input("Due Date", value=default_due, key="due_date")
    biller_name = st.text_input("Biller Name", value="Mr. Manish Dugar", key="biller_name")

    st.markdown(f"**Invoice Number:** `{bill_no}`")

    # ── NEW: Client Info Section ──
    st.markdown("### Client Information")
    client_address = st.text_input(
        "Client Address (max 65 characters)",
        value="",
        max_chars=65,
        placeholder="Enter client address",
        key="client_address",
    )
    client_phone = st.text_input(
        "Client Phone Number",
        value="",
        placeholder="Enter client phone number",
        key="client_phone",
    )
    client_email = st.text_input(
        "Client Email",
        value="",
        placeholder="Enter client email",
        key="client_email",
    )
    client_bill_to = st.text_input(
        "Client Bill To",
        value="",
        placeholder="Enter client billing name",
        key="client_bill_to",
    )

    # ── NEW: Payment Method Section ──
    st.markdown("### Payment Method")
    payment_method = st.radio(
        "Select Payment Method:",
        options=["Cash", "NEFT / IMPS", "UPI", "Cheque"],
        index=0,  # Default to "Cash"
        key="payment_method",
    )

    # Initialize session-state rows
    if "rows" not in st.session_state:
        st.session_state.rows = [
            {"No.": "", "Item Description": "", "Weight": "", "Rate (₹)": "", "Amount (₹)": ""}
        ]

    rows = st.session_state.rows

    st.markdown("---")

    # Always use vertical card per line item, with label above each input (responsive & clear)
    for idx in range(len(rows)):
        with st.container():
            st.markdown(f"**Line Item {idx+1}**")
            c1, c2 = st.columns([1, 5], gap="small")
            with c1:
                st.markdown("No. (positive integer)")
            with c2:
                no_val = st.text_input(
                    label="No. (positive integer)",
                    value=rows[idx]["No."],
                    placeholder="1",
                    key=f"No_{idx}",
                    label_visibility="collapsed",
                )
            c1, c2 = st.columns([1, 5], gap="small")
            with c1:
                st.markdown("Item Description (required)")
            with c2:
                desc_val = st.text_input(
                    label="Item Description (required)",
                    value=rows[idx]["Item Description"],
                    placeholder="Diamond Ring, Necklace…",
                    key=f"Desc_{idx}",
                    label_visibility="collapsed",
                )
            c1, c2 = st.columns([1, 5], gap="small")
            with c1:
                st.markdown("Weight (gm) (non-negative)")
            with c2:
                weight_val = st.text_input(
                    label="Weight (gm) (non-negative)",
                    value=rows[idx]["Weight"],
                    placeholder="e.g. 1.25",
                    key=f"Weight_{idx}",
                    label_visibility="collapsed",
                )
            c1, c2 = st.columns([1, 5], gap="small")
            with c1:
                st.markdown("Rate (₹) (non-negative)")
            with c2:
                rate_val = st.text_input(
                    label="Rate (₹) (non-negative)",
                    value=rows[idx]["Rate (₹)"],
                    placeholder="e.g. 45000",
                    key=f"Rate_{idx}",
                    label_visibility="collapsed",
                )
            # --- Auto-calculate Amount ---
            try:
                w = float(weight_val)
                r = float(rate_val)
                auto_amount = f"{w * r:.2f}"
            except Exception:
                auto_amount = ""
            c1, c2 = st.columns([1, 5], gap="small")
            with c1:
                st.markdown("Amount (₹) (auto, editable)")
            with c2:
                amount_val = st.text_input(
                    label="Amount (₹) (auto, editable)",
                    value=rows[idx]["Amount (₹)"] if rows[idx]["Amount (₹)"] else auto_amount,
                    placeholder=auto_amount,
                    key=f"Amount_{idx}",
                    label_visibility="collapsed",
                )
            # Save back to session state
            st.session_state.rows[idx]["No."] = no_val
            st.session_state.rows[idx]["Item Description"] = desc_val
            st.session_state.rows[idx]["Weight"] = weight_val
            st.session_state.rows[idx]["Rate (₹)"] = rate_val
            st.session_state.rows[idx]["Amount (₹)"] = amount_val

    st.markdown("---")

    col_a, col_b = st.columns([1, 1])
    with col_a:
        if st.button("➕ Add Another Row"):
            st.session_state.rows.append(
                {
                    "No.": "",
                    "Item Description": "",
                    "Weight": "",
                    "Rate (₹)": "",
                    "Amount (₹)": "",
                }
            )
            st.rerun()

    with col_b:
        generate_button = st.button("🖨️ Generate Invoice")

    # When “Generate Invoice” is clicked:
    if generate_button:
        # 1) Filter out fully blank rows
        filtered_rows = []
        for r in st.session_state.rows:
            if any(str(v).strip() != "" for v in r.values()):
                filtered_rows.append(r.copy())

        if not filtered_rows:
            st.error("No data entered. Please fill at least one line item before generating the invoice.")
            return

        # 2) Stricter validation
        errors = []
        validated = []
        for idx, r in enumerate(filtered_rows, start=1):
            row_errs = []
            # Validate No. (positive integer)
            try:
                no_int = int(parse_number(r["No."]))
                if no_int <= 0:
                    row_errs.append("No. must be a positive integer.")
            except Exception:
                row_errs.append("No. must be a positive integer.")

            # Validate Item Description (non-empty)
            if not str(r["Item Description"]).strip():
                row_errs.append("Item Description cannot be empty.")

            # Validate Weight (non-negative float)
            try:
                w = parse_number(r["Weight"])
                if w < 0:
                    row_errs.append("Weight must be non-negative.")
            except Exception:
                row_errs.append("Weight must be a number (e.g. 1.25).")

            # Validate Rate (non-negative float)
            try:
                rt = parse_number(r["Rate (₹)"])
                if rt < 0:
                    row_errs.append("Rate (₹) must be non-negative.")
            except Exception:
                row_errs.append("Rate (₹) must be a number (e.g. 45000).")

            # Validate Amount (non-negative float), allow blank and auto-calc
            amt_val = r["Amount (₹)"]
            amt_auto = False
            try:
                if str(amt_val).strip() == "":
                    # Try to auto-calculate
                    amt = parse_number(r["Weight"]) * parse_number(r["Rate (₹)"])
                    amt_auto = True
                else:
                    amt = parse_number(amt_val)
                if amt < 0:
                    row_errs.append("Amount (₹) must be non-negative.")
            except Exception:
                row_errs.append("Amount (₹) must be a number (e.g. 56250) or left blank for auto-calc.")

            if row_errs:
                errors.append(f"Row {idx}: " + "; ".join(row_errs))
            else:
                validated.append(
                    {
                        "No.": str(int(parse_number(r["No."]))),
                        "Item Description": r["Item Description"].strip(),
                        "Weight": f"{parse_number(r['Weight']):.2f}",
                        "Rate (₹)": f"{parse_number(r['Rate (₹)']):.2f}",
                        "Amount (₹)": f"{amt:.2f}",
                    }
                )

        if errors:
            st.error("Please fix these errors before generating the invoice:")
            for e in errors:
                st.write(f"- {e}")
            return

        # 3) Sort by No.
        validated.sort(key=lambda x: int(x["No."]))

        # 4) Generate PPTX
        bill_info = {
            "Bill No": bill_no,
            "Bill Date": bill_date.strftime("%d-%m-%Y"),
            "Due Date": due_date.strftime("%d-%m-%Y"),
            "Biller Name": biller_name,
            "Client Address": client_address,
            "Client Phone Number": client_phone,
            "Client Email": client_email,
            "Client Bill To": client_bill_to,
        }
        try:
            pptx_bytes = generate_filled_invoice(validated, "invoice_template.pptx", bill_info, payment_method)
        except Exception as e:
            st.error(f"Unexpected error during PPT generation: {e}")
            return

        # Generate a dynamic filename
        filename = f"RishabGems_{bill_no}_{client_bill_to}_{client_phone}_{today.strftime('%Y%m%d')}.pptx"
        filename = re.sub(r'[\\/*?:"<>|]', "", filename)  # Remove invalid characters for filenames

        # 5) Download button for PPTX
        if pptx_bytes:
            st.success("✅ Invoice generated successfully.")
            st.download_button(
                label="📥 Download Filled PPTX",
                data=pptx_bytes,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )


if __name__ == "__main__":
    main()
